"""Fill the official UiPath AgentHack template (7 slides) from a content.json file.

Theme is preserved from the template master (Segoe UI). Only font size/bold/name are set.
See ../THEME.md for the placeholder schema and ../SKILL.md for the workflow.

Usage:
  python build_deck.py --content content.json --out out.pptx [--template t.pptx] [--arch-image a.png]
"""
import argparse
import json
import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

# Candidate template locations, in priority order (after --template).
TEMPLATE_CANDIDATES = [
    Path(r"C:/Users/DanielaRosenstein/my-ds-project/Submission deck.pptx"),
    Path("submission/deck/Submission-Template-Deal-Desk-Agent.pptx"),
]
DEFAULT_ARCH = Path("submission/code-graph/deal-desk-code-graph.png")

# Friendly content key -> (placeholder substring, font size, bold_first, exact_match)
SCHEMA = {
    "title":       ("Presentation title goes here", 40, True, False),
    "team_footer": ("Team name", 16, False, False),
    "problem":     ("What real-world problem", 13, False, False),
    "solution":    ("Brief summary of the solution", 12.5, False, False),
    "benefits":    ("What does this agent actually achieve", 12, False, False),
    "tech":        ("Details", 11.5, True, True),
    "architecture":("This slide is optional", 10, False, False),
    "demo_header": ("Miscellaneous", 24, True, True),
    "demo_body":   ("If you need extra slides", 12, False, False),
    "closing":     ("Closing message", 24, True, False),
}
NAME_PH, ROLE_PH = "Jane Doe", "Job title @Company"


def shape_text(shape):
    if not getattr(shape, "has_text_frame", False):
        return ""
    return "\n".join(p.text for p in shape.text_frame.paragraphs).strip()


def set_block(shape, lines, size=18, bold_first=True, align=PP_ALIGN.LEFT):
    """Replace all text in a shape, preserving Segoe UI; blank strings = blank lines."""
    tf = shape.text_frame
    tf.clear()
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.alignment = align
        if p.runs:
            r = p.runs[0]
            r.font.size = Pt(size)
            r.font.bold = bool(bold_first and i == 0)
            r.font.name = "Segoe UI"


def matches(text, needle, exact):
    return text == needle if exact else (needle in text)


def resolve_template(arg_template):
    if arg_template:
        p = Path(arg_template)
        if not p.exists():
            sys.exit(f"ERROR: --template not found: {p}")
        return p
    for c in TEMPLATE_CANDIDATES:
        if c.exists():
            return c
    sys.exit("ERROR: no template found. Pass --template <official AgentHack template.pptx>. "
             "See .cursor/skills/presentation/THEME.md.")


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("--content", required=True, help="content.json path")
    ap.add_argument("--out", required=True, help="output .pptx path")
    ap.add_argument("--template", help="base template .pptx (official AgentHack deck)")
    ap.add_argument("--arch-image", help="architecture image for slide 5")
    args = ap.parse_args()

    content = json.loads(Path(args.content).read_text(encoding="utf-8"))
    template = resolve_template(args.template)
    prs = Presentation(str(template))

    # Generic single-shape placeholders.
    for slide in prs.slides:
        for sh in slide.shapes:
            t = shape_text(sh)
            if not t:
                continue
            for key, (needle, size, bold, exact) in SCHEMA.items():
                if key in content and matches(t, needle, exact):
                    set_block(sh, content[key], size=size, bold_first=bold)
                    break

    # Team slide: repeated name/role slots filled positionally.
    names = content.get("team_names", [])
    roles = content.get("team_roles", [])
    if names or roles:
        for slide in prs.slides:
            name_slots = [sh for sh in slide.shapes if shape_text(sh) == NAME_PH]
            role_slots = [sh for sh in slide.shapes if ROLE_PH in shape_text(sh)]
            if not (name_slots or role_slots):
                continue
            for sh, v in zip(name_slots, names):
                set_block(sh, [v], size=18)
            for sh, v in zip(role_slots, roles):
                set_block(sh, [v], size=13, bold_first=False)
            break

    # Architecture image on the slide carrying the architecture placeholder.
    arch = Path(args.arch_image) if args.arch_image else DEFAULT_ARCH
    if "architecture" in content and arch.exists():
        for slide in prs.slides:
            if any("This slide is optional" in shape_text(sh) for sh in slide.shapes):
                slide.shapes.add_picture(str(arch), Inches(0.95), Inches(4.0),
                                         width=Inches(11.4), height=Inches(2.6))
                break

    out = Path(args.out)
    out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out))
    print(f"Saved: {out}")
    print(f"  template={template.name}  slides={len(prs.slides)}  size={out.stat().st_size/1e6:.2f} MB")


if __name__ == "__main__":
    main()
