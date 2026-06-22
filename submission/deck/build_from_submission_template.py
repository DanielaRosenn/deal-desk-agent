"""
Fill the official UiPath AgentHack submission template (7 slides) with Deal Desk Agent content.

Narrative: "Every sales org has a Deal Desk. Ours is an AI agent."
Three-channel HITL: Outlook Adaptive Card + Action Center + Slack (all three verified live Jun 22 2026).

Accurate architecture (v1.3.16 Robot / v1.1.9 BPMN):
- Approval delivery: Outlook Adaptive Card (originator 61fed71d, Action.Http buttons, HTML fallback)
- HITL: UiPath Action Center external task (Deal Desk catalog, UiPath.Persistence.Activities)
- Slack: Block Kit DM to approver + requester
- Robot: WaitDecision UiPath RPA v1.3.16 — handles all three channels, suspends via Persistence
- No BPMN polling loop — robot blocks internally and returns decision directly
- Salesforce: live IS + AgentHub MCP sidecar

Run:  python build_from_submission_template.py
"""

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

TEMPLATE = Path(r"c:\Users\DanielaRosenstein\my-ds-project\Submission deck.pptx")
OUT = Path(r"c:\Users\DanielaRosenstein\OneDrive - Cato Networks\Documents\UiPath\RenewalPrice_Flow\submission\deck\Submission-Template-Deal-Desk-Agent.pptx")
ARCH_IMG = Path(r"c:\Users\DanielaRosenstein\OneDrive - Cato Networks\Documents\UiPath\RenewalPrice_Flow\submission\code-graph\deal-desk-code-graph.png")


def shape_text(shape):
    if not getattr(shape, "has_text_frame", False):
        return ""
    return "\n".join(p.text for p in shape.text_frame.paragraphs).strip()


def set_block(shape, lines, size=18, bold_first=True, align=PP_ALIGN.LEFT):
    """Replace all text in a shape with new lines, preserving the first run's font."""
    tf = shape.text_frame
    tf.clear()
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.alignment = align
        if p.runs:
            r = p.runs[0]
            r.font.size = Pt(size)
            r.font.bold = bool(bold_first and i == 0)
            r.font.name = "Segoe UI"


prs = Presentation(str(TEMPLATE))

# ---------------------------------------------------------------------------
# Slide 1 — Title
# ---------------------------------------------------------------------------
s1 = prs.slides[0]
for sh in s1.shapes:
    if "Presentation title goes here" in shape_text(sh):
        set_block(sh, [
            "Deal Desk Agent",
            "Your AI Deal Desk Analyst",
        ], size=40)

# ---------------------------------------------------------------------------
# Slide 2 — Team / contributors
# ---------------------------------------------------------------------------
s2 = prs.slides[1]
name_slots, role_slots = [], []
team_slot = None
for sh in s2.shapes:
    t = shape_text(sh)
    if t == "Jane Doe":
        name_slots.append(sh)
    elif "Job title @Company" in t:
        role_slots.append(sh)
    elif "Team name" in t:
        team_slot = sh

names = ["Daniela Rosenstein", "LangGraph coded agent", "Maestro BPMN", "WaitDecision Robot"]
roles = [
    "Builder @ Cato Networks",
    "plan / render / process_response (Python LangGraph)",
    "orchestration, gateway, sequential loop, polling",
    "RPA: POSTs to HITL portal, sends Outlook email",
]
for sh, v in zip(name_slots, names):
    set_block(sh, [v], size=18)
for sh, v in zip(role_slots, roles):
    set_block(sh, [v], size=13, bold_first=False)
if team_slot:
    set_block(team_slot, [
        "Track 2: UiPath Maestro BPMN",
        "Coding-agent bonus: Cursor + Claude via UiPath for Coding Agents",
    ], size=16)

# ---------------------------------------------------------------------------
# Slide 3 — Problem + Solution
# ---------------------------------------------------------------------------
s3 = prs.slides[2]
for sh in s3.shapes:
    t = shape_text(sh)
    if "What real-world problem" in t:
        set_block(sh, [
            "Every sales org has a Deal Desk: the person who reviews discount requests,",
            "emails managers up the chain, waits for replies, and updates the CRM.",
            "It is entirely manual, inbox-driven, and slower than the sales team needs.",
            "",
            "The Opportunity Owner (AE) submits a request and then waits:",
            "  - No SLA: deals stall in inboxes for days.",
            "  - No consistency: routing thresholds drift between managers.",
            "  - No audit trail: decisions buried in email threads.",
        ], size=13, bold_first=False)
    elif "Brief summary of the solution" in t:
        set_block(sh, [
            "The Deal Desk Agent does the Deal Desk job:",
            "",
            "  1. Queries live Salesforce data via AgentHub MCP sidecar.",
            "  2. Resolves approver chain from HiBob org chart (reportsTo hierarchy).",
            "  3. Applies risk scoring + LLM reasoning to draft a recommendation.",
            "  4. WaitDecision Robot: Action Center task + Outlook Adaptive Card + Slack DM",
            "     — three channels simultaneously. Robot suspends via UiPath Persistence.",
            "  5. Approver decides in Outlook (Adaptive Card button) or Action Center.",
            "  6. Robot resumes, returns decision to BPMN. Rejection short-circuits chain.",
            "  7. Slack + Outlook summary to AE. Data Fabric audit record.",
        ], size=12.5, bold_first=False)

# ---------------------------------------------------------------------------
# Slide 4 — Benefits + Technologies
# ---------------------------------------------------------------------------
s4 = prs.slides[3]
for sh in s4.shapes:
    t = shape_text(sh)
    if "What does this agent actually achieve" in t:
        set_block(sh, [
            "For the Opportunity Owner (AE / requester):",
            "  - No more chasing managers: agent routes, robot sends, BPMN tracks.",
            "  - Slack DM + Outlook summary with full decision trail on completion.",
            "  - Low-risk deals (<=5% discount / <=$50k) clear with ZERO human touch.",
            "",
            "For the business:",
            "  - Consistent policy applied identically to every deal (risk score formula).",
            "  - SLA: BPMN polls every 30s; timer fires reminder via Robot if silent.",
            "  - Rejection short-circuits chain immediately (no wasted approver time).",
            "  - Audit: Data Fabric logs agent recommendation vs. each human decision.",
        ], size=12, bold_first=False)
    elif t.strip() == "Details":
        set_block(sh, [
            "Technologies (all verified live Jun 22 2026):",
            "  UiPath Coded Agents (Python + LangGraph): plan, render, process_response",
            "  UiPath Maestro BPMN 2.0: sequential multi-instance loop, completionCondition",
            "  WaitDecision Robot v1.3.16 (UiPath RPA + Persistence.Activities): all channels",
            "  Action Center: external task in Deal Desk catalog, Approve/Reject buttons",
            "  Outlook Adaptive Card: originator 61fed71d, Action.Http buttons, HTML fallback",
            "  Slack Block Kit: DM to approver + requester via Slack Bot Token",
            "  AgentHub MCP Server: Salesforce DealDesk MCP (live IS, read-only SOQL)",
            "  Data Fabric: ApprovalAudit entity (agent rec vs. each human decision)",
            "  Solution deployed to Shared/DealDeskApprovalGlobal on Automation Cloud",
        ], size=11.5)

# ---------------------------------------------------------------------------
# Slide 5 — Architecture
# ---------------------------------------------------------------------------
s5 = prs.slides[4]
for sh in s5.shapes:
    t = shape_text(sh)
    if "This slide is optional" in t:
        set_block(sh, [
            "Salesforce event -> Maestro BPMN -> plan agent:",
            "  AgentHub MCP sidecar queries live Salesforce Opportunity + Account data",
            "  HiBob org chart resolution -> builds approver chain (manager->director->vp->cfo->cro)",
            "  risk_score formula + LLM reasoning -> RoutingPlan (disposition + rationale + chain)",
            "",
            "Disposition gateway:",
            "  auto_approve  ->  Slack + Outlook to AE + Data Fabric audit  (zero human touch)",
            "  exception     ->  immediate human review",
            "  needs_approval  ->  sequential approver loop (per-approver):",
            "    render agent builds Adaptive Card payload (deal context + LLM recommendation)",
            "    WaitDecision Robot v1.3.16:",
            "      Action Center: CreateExternalTask ('Deal Desk' catalog)",
            "      Outlook: Adaptive Card with Approve / Reject / Request Info buttons",
            "      Slack: Block Kit DM to approver",
            "      Suspends via UiPath.Persistence.Activities (no BPMN polling)",
            "    Approver decides via Adaptive Card OR Action Center -> robot resumes",
            "    Rejection short-circuits chain (BPMN completionCondition)",
            "",
            "On completion: Slack DM + Outlook summary to AE -> Data Fabric ApprovalAudit",
        ], size=10.0, bold_first=False)

if ARCH_IMG.exists():
    s5.shapes.add_picture(str(ARCH_IMG), Inches(0.95), Inches(4.0), width=Inches(11.4), height=Inches(2.6))

# ---------------------------------------------------------------------------
# Slide 6 — Demo plan
# ---------------------------------------------------------------------------
s6 = prs.slides[5]
for sh in s6.shapes:
    t = shape_text(sh)
    if t.strip() == "Miscellaneous":
        set_block(sh, ["Three demo scenarios"], size=24)
    elif "If you need extra slides" in t:
        set_block(sh, [
            "Live test results — v1.3.12, Jun 22 2026 (all channels verified):",
            "  TEST-001: GlobalTech $120k/18% -> Agent: Approve (within policy)",
            "  TEST-002: MegaCorp $350k/32%   -> Agent: Approve with conditions (escalation)",
            "  TEST-003: StartupXYZ $28k/45%  -> Agent: Reject (margin impact too high)",
            "  TEST-004: Pharma $780k/22%     -> Agent: Approve (enterprise, Finance pre-approved)",
            "  Outlook Adaptive Card RECEIVED  |  Slack DM RECEIVED  |  Action Center WORKING",
            "",
            "Scenario 1 — Auto-approve  (4% / $30k)",
            "  Disposition: auto_approve. AE notified instantly. Zero human touches.",
            "",
            "Scenario 2 — One approver  (8% / $40k)",
            "  Robot sends Adaptive Card + Slack to Maya Stone (Manager).",
            "  Maya approves via Outlook card. AE gets Slack DM + Outlook summary.",
            "",
            "Scenario 3 — Full chain  (30% / $1.2M)",
            "  Four sequential approvers. Each gets Adaptive Card + Slack.",
            "  Rejection at any step short-circuits. Full Data Fabric audit trail.",
        ], size=12.0, bold_first=False)

# ---------------------------------------------------------------------------
# Slide 7 — Closing
# ---------------------------------------------------------------------------
s7 = prs.slides[6]
for sh in s7.shapes:
    t = shape_text(sh)
    if "Closing message" in t:
        set_block(sh, [
            "Every Deal Desk should work like this.",
            "Collect. Reason. Route. Sign-off. Audit. Automatically.",
            "",
            "Track 2: Maestro BPMN  |  Coding-agent bonus  |  Runs on Automation Cloud",
        ], size=24)

OUT.parent.mkdir(parents=True, exist_ok=True)
prs.save(str(OUT))
print(f"Saved: {OUT}")
print(f"  slides={len(prs.slides)}  size={OUT.stat().st_size / 1e6:.2f} MB")
