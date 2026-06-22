"""Build a fully editable PowerPoint deck (native shapes/text, no full-slide images)."""
from __future__ import annotations

import os

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

HERE = os.path.dirname(os.path.abspath(__file__))
GRAPH_PNG = os.path.join(HERE, "..", "code-graph", "deal-desk-code-graph.png")

INK = RGBColor(0x11, 0x16, 0x1C)
PANEL = RGBColor(0x1C, 0x23, 0x2B)
PANEL2 = RGBColor(0x16, 0x1C, 0x23)
ORANGE = RGBColor(0xFA, 0x46, 0x16)
BLUE = RGBColor(0x2E, 0xA8, 0xFF)
CYAN = RGBColor(0x22, 0xD3, 0xEE)
AMBER = RGBColor(0xFF, 0xB0, 0x20)
TEAL = RGBColor(0x2B, 0xD9, 0xA6)
VIOLET = RGBColor(0xB7, 0x94, 0xF6)
WHITE = RGBColor(0xF5, 0xF8, 0xFB)
SLATE = RGBColor(0xC9, 0xD3, 0xDD)
MUTED = RGBColor(0x8A, 0x97, 0xA4)

NOTES = [
    "UiPath's 2026 challenge: bridge the gap between a prototype on a laptop and software running in production.",
    "Manual AS-IS is inbox ping-pong: slow, inconsistent, unauditable.",
    "Core design: agent owns judgment; Maestro owns lifecycle.",
    "Three LangGraph entry points. plan produces approver chain as data.",
    "One BPMN process with multi-instance + SLA timer + audit.",
    "Narrate trigger -> decide -> route -> human -> resume -> audit loop.",
    "One process, three behaviors: auto-clear, 1-tier, multi-tier.",
    "Built with coding agents (Cursor + Claude via UiPath for Coding Agents).",
    "Governance checklist: human gate, SLA, audit, guardrails.",
    "Everything runs on Automation Cloud; packaged as .uipx.",
    "ROI: deflection, speed, consistency, auditability; portable pattern.",
    "Replace links and close with interactive page CTA.",
]


def set_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, x, y, w, h, text, size=18, color=SLATE, bold=False, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    p.level = 0
    r = p.runs[0]
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    r.font.name = "Segoe UI"
    return box


def add_card(slide, x, y, w, h, title, lines, accent=ORANGE):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid(); shp.fill.fore_color.rgb = PANEL
    shp.line.color.rgb = accent; shp.line.width = Pt(1.4)

    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(0.07))
    bar.fill.solid(); bar.fill.fore_color.rgb = accent; bar.line.fill.background()

    add_textbox(slide, x + 0.18, y + 0.13, w - 0.35, 0.32, title, size=14, color=WHITE, bold=True)
    yline = y + 0.5
    for line in lines:
        add_textbox(slide, x + 0.18, yline, w - 0.35, 0.3, line, size=11.5, color=SLATE)
        yline += 0.28


def add_chrome(slide, idx, kicker="TRACK 2 · UIPATH MAESTRO BPMN"):
    mark = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.45), Inches(0.22), Inches(0.13), Inches(0.13))
    mark.fill.solid(); mark.fill.fore_color.rgb = ORANGE; mark.line.fill.background()
    add_textbox(slide, 0.62, 0.17, 4.8, 0.26, "UiPath AgentHack 2026", size=14, color=WHITE, bold=True)
    add_textbox(slide, 8.2, 0.17, 4.9, 0.26, kicker, size=11.5, color=ORANGE, bold=True, align=PP_ALIGN.RIGHT)

    f = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(7.36), Inches(13.333), Inches(0.04))
    f.fill.solid(); f.fill.fore_color.rgb = RGBColor(0x22, 0x2A, 0x33); f.line.fill.background()
    add_textbox(slide, 0.45, 7.08, 7.0, 0.2, "Deal Desk Agent - from laptop to production, governed", size=9, color=MUTED)

    tag = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(12.2), Inches(6.92), Inches(0.95), Inches(0.32))
    tag.fill.solid(); tag.fill.fore_color.rgb = ORANGE; tag.line.fill.background()
    add_textbox(slide, 12.2, 6.95, 0.95, 0.2, f"{idx:02d}/12", size=10, color=INK, bold=True, align=PP_ALIGN.CENTER)


def add_title(slide, text):
    add_textbox(slide, 0.45, 0.95, 10.2, 1.1, text, size=34, color=WHITE, bold=True)


def add_bullets(slide, x, y, items, size=15, marker=ORANGE):
    yy = y
    for item in items:
        dot = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(yy + 0.1), Inches(0.08), Inches(0.08))
        dot.fill.solid(); dot.fill.fore_color.rgb = marker; dot.line.fill.background()
        add_textbox(slide, x + 0.2, yy, 12.2 - x, 0.34, item, size=size, color=SLATE)
        yy += 0.44


def set_notes(slide, note):
    slide.notes_slide.notes_text_frame.text = note


def build():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    s = prs.slides.add_slide(blank); set_bg(s, INK); add_chrome(s, 1); add_title(s, "Deal Desk Agent")
    add_textbox(s, 0.45, 2.08, 10, 0.5, "From laptop to production - governed.", size=24, color=ORANGE, bold=True)
    add_textbox(s, 0.45, 2.66, 10, 0.3, "An autonomous pricing-approval agent on UiPath Maestro BPMN.", size=15, color=SLATE)
    add_card(s, 0.45, 5.45, 3.95, 0.68, "Track 2: Maestro BPMN", [], ORANGE)
    add_card(s, 4.7, 5.45, 3.95, 0.68, "Built with coding agents", [], BLUE)
    add_card(s, 8.95, 5.45, 3.95, 0.68, "Runs on Automation Cloud", [], TEAL)
    set_notes(s, NOTES[0])

    s = prs.slides.add_slide(blank); set_bg(s, INK); add_chrome(s, 2)
    add_title(s, "Pricing approvals: 100% human-touched, mostly for nothing")
    add_bullets(s, 0.5, 2.2, [
        "Every discount / renewal change is routed to a human - even the trivial ones",
        "Approval thresholds drift between managers",
        "No SLA and no auditable decision trail",
    ], size=16)
    add_card(s, 0.45, 5.05, 4.1, 1.15, "Slow", ["inbox ping-pong"], ORANGE)
    add_card(s, 4.65, 5.05, 4.1, 1.15, "Inconsistent", ["tribal thresholds"], AMBER)
    add_card(s, 8.85, 5.05, 4.1, 1.15, "Invisible", ["no audit trail"], VIOLET)
    set_notes(s, NOTES[1])

    s = prs.slides.add_slide(blank); set_bg(s, INK); add_chrome(s, 3)
    add_title(s, "An agent for judgment. Maestro for the lifecycle.")
    add_card(s, 0.45, 2.3, 6.1, 3.45, "AGENT (LangGraph)", [
        "Decides what should happen", "Disposition + risk score",
        "Recommendation + rationale", "Sizes approver chain as data",
    ], BLUE)
    add_card(s, 6.8, 2.3, 6.1, 3.45, "MAESTRO BPMN", [
        "Runs how it happens", "Cards, human gate, SLA timer",
        "Iterates agent chain", "Audit on every path",
    ], ORANGE)
    set_notes(s, NOTES[2])

    s = prs.slides.add_slide(blank); set_bg(s, INK); add_chrome(s, 4)
    add_title(s, "The agent: Python LangGraph coded agent")
    add_card(s, 0.45, 2.2, 12.45, 2.15, "Entry points", [
        "plan: score, set disposition, draft recommendation, size approver chain",
        "render: generate contextual Outlook Adaptive Card per approver",
        "process_response: interpret human reply and pick next action",
    ], BLUE)
    logic = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.45), Inches(4.62), Inches(12.45), Inches(1.55))
    logic.fill.solid(); logic.fill.fore_color.rgb = PANEL2; logic.line.color.rgb = AMBER
    add_textbox(s, 0.7, 4.78, 12.0, 0.3, "Decision logic (from code)", size=12.5, color=AMBER, bold=True)
    add_textbox(s, 0.7, 5.08, 12.0, 0.85,
                "risk = min(1.0, 0.5*(discount/25%) + 0.5*(value/$500k))\n"
                "auto_clear if <=5% and <=$50k; >25% adds Director+VP; >$500k adds CFO; invalid payload -> exception",
                size=11, color=SLATE)
    set_notes(s, NOTES[3])

    s = prs.slides.add_slide(blank); set_bg(s, INK); add_chrome(s, 5)
    add_title(s, "One Maestro BPMN process - right actor per step")
    add_bullets(s, 0.5, 2.1, [
        "Start (opportunity change) -> agent plan",
        "Disposition gateway: auto_clear / exception close out immediately",
        "needs_approval -> multi-instance subprocess over approver chain",
        "Boundary SLA timer -> reminder + escalation; Data Fabric audit",
    ], size=15)
    set_notes(s, NOTES[4])

    s = prs.slides.add_slide(blank); set_bg(s, INK); add_chrome(s, 6)
    add_title(s, "The whole flow - explore it")
    add_textbox(s, 0.5, 2.0, 12.1, 0.3,
                "Salesforce + HiBob -> LangGraph -> Maestro BPMN -> Outlook / Action Center / Azure bridge / Data Fabric",
                size=12.5, color=SLATE)
    if os.path.exists(GRAPH_PNG):
        s.shapes.add_picture(GRAPH_PNG, Inches(1.05), Inches(2.35), Inches(11.2), Inches(4.15))
    set_notes(s, NOTES[5])

    s = prs.slides.add_slide(blank); set_bg(s, INK); add_chrome(s, 7)
    add_title(s, "Three paths, one process")
    add_card(s, 0.45, 2.3, 4.1, 3.4, "1) Auto-clear", ["<=5% / <=$50k", "Zero cards", "Logged + closed"], TEAL)
    add_card(s, 4.65, 2.3, 4.1, 3.4, "2) One manager", ["8% / $40k", "One card", "Approve -> approved"], BLUE)
    add_card(s, 8.85, 2.3, 4.1, 3.4, "3) Multi-tier", ["30% / $1.2M", "Manager -> Director -> VP -> CFO", "Sequential cards"], ORANGE)
    set_notes(s, NOTES[6])

    s = prs.slides.add_slide(blank); set_bg(s, INK); add_chrome(s, 8, "BONUS · UIPATH FOR CODING AGENTS")
    add_title(s, "Built by coding agents")
    add_textbox(s, 0.5, 1.92, 10.6, 0.35, "Natural language -> coding agents -> governed cloud process", size=16, color=ORANGE, bold=True)
    add_card(s, 0.45, 2.5, 3.1, 1.35, "Prompts", ["Laptop intent"], MUTED)
    add_card(s, 3.75, 2.5, 3.1, 1.35, "Cursor + Claude", ["uip skills + CLI"], BLUE)
    add_card(s, 7.05, 2.5, 2.8, 1.35, "Generated", ["Agent + BPMN + DF"], VIOLET)
    add_card(s, 10.05, 2.5, 2.85, 1.35, "Automation Cloud", ["Published + deployed"], TEAL)
    add_bullets(s, 0.5, 4.1, [
        "Coding agents generated/edited the LangGraph agent, BPMN, bindings, schema exports, Data Fabric entity",
        "Low-code Maestro + pro-code agent combined",
        "Embodies the laptop -> production bridge",
    ], size=13.5)
    set_notes(s, NOTES[7])

    s = prs.slides.add_slide(blank); set_bg(s, INK); add_chrome(s, 9)
    add_title(s, "Governance & human-in-the-loop")
    add_card(s, 0.45, 2.35, 3.05, 3.35, "Human gate", ["Action Center task", "Azure bridge resumes instance"], VIOLET)
    add_card(s, 3.75, 2.35, 3.05, 3.35, "Time governance", ["Boundary SLA timer", "Reminder + escalation"], AMBER)
    add_card(s, 7.05, 2.35, 3.05, 3.35, "Audit", ["Data Fabric logs", "agent rec vs human decision"], TEAL)
    add_card(s, 10.35, 2.35, 2.95, 3.35, "Guardrails", ["Invalid payload -> exception", "Empty chain -> manager fallback"], ORANGE)
    set_notes(s, NOTES[8])

    s = prs.slides.add_slide(blank); set_bg(s, INK); add_chrome(s, 10)
    add_title(s, "Runs on Automation Cloud")
    add_bullets(s, 0.5, 2.1, [
        "Coded agent published to tenant feed; 3 entry points deployed as runnable processes",
        "BPMN validated; live Outlook connection bound",
        "Full solution packaged as deployable .uipx (agent + BPMN + Data Fabric ApprovalAudit)",
        "Unit tests passing across disposition, auto-clear, multi-tier, schema, guardrails, rendering",
    ], size=14.5)
    add_textbox(s, 0.5, 6.2, 12.0, 0.35, "UiPath is the orchestration + governance layer that ties it together.", size=14.5, color=ORANGE, bold=True)
    set_notes(s, NOTES[9])

    s = prs.slides.add_slide(blank); set_bg(s, INK); add_chrome(s, 11)
    add_title(s, "Business impact & reach")
    add_card(s, 0.45, 2.35, 3.05, 3.3, "Deflection", ["Low-risk deals clear", "with zero human touch"], TEAL)
    add_card(s, 3.75, 2.35, 3.05, 3.3, "Speed + consistency", ["SLA + routing replace", "inbox ping-pong"], BLUE)
    add_card(s, 7.05, 2.35, 3.05, 3.3, "Auditability", ["Every decision logged", "for compliance + tuning"], VIOLET)
    add_card(s, 10.35, 2.35, 2.95, 3.3, "Portable", ["Invoice / expense / onboarding", "same pattern"], ORANGE)
    set_notes(s, NOTES[10])

    s = prs.slides.add_slide(blank); set_bg(s, INK); add_chrome(s, 12)
    add_title(s, "Decide. Draft. Route. Audit.")
    add_textbox(s, 0.45, 2.08, 11.5, 0.45, "Autonomously - and shipped to production.", size=23, color=ORANGE, bold=True)
    add_textbox(s, 0.45, 2.66, 12.2, 0.35, "Agent owns judgment · Maestro owns lifecycle · UiPath owns governance", size=13.5, color=SLATE)
    add_card(s, 0.45, 4.0, 4.05, 1.2, "DevPost", ["<link>"], BLUE)
    add_card(s, 4.65, 4.0, 4.05, 1.2, "Demo video", ["<link>"], BLUE)
    add_card(s, 8.85, 4.0, 4.05, 1.2, "Code", ["<link>"], BLUE)
    set_notes(s, NOTES[11])

    out = os.path.join(HERE, "Deal-Desk-Agent-native-editable.pptx")
    prs.save(out)
    print(f"{out} ({len(prs.slides)} slides, {os.path.getsize(out)/1e6:.2f} MB)")


if __name__ == "__main__":
    build()
